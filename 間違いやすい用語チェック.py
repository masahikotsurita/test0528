import os
import sys
import datetime
import configparser
import tkinter as tk
from tkinter import ttk, messagebox

if os.name == 'nt':
    try:
        import ctypes
        ctypes.windll.user32.ShowWindow(ctypes.windll.kernel32.GetConsoleWindow(), 0)
    except Exception:
        pass

# iniファイルからkeywordsとreplacementsを読み込む
config = configparser.ConfigParser()

# 実行ファイルの配置ディレクトリを基準とする
if getattr(sys, 'frozen', False):
    base_path = os.path.dirname(sys.executable)  # PyInstaller実行環境
else:
    base_path = os.path.dirname(__file__)

# PyInstallerの --add-data でバンドルされた ini も探索する
ini_candidates = [os.path.join(base_path, '間違いやすい用語チェック.ini')]
if getattr(sys, '_MEIPASS', None):
    ini_candidates.append(os.path.join(sys._MEIPASS, '間違いやすい用語チェック.ini'))

ini_path = None
for p in ini_candidates:
    if os.path.exists(p):
        ini_path = p
        break
if ini_path is None:
    ini_path = ini_candidates[0]

def load_replacements():
    """Load keywords and replacements from the selected section."""
    config.read(ini_path, encoding="utf-8")

    # Determine the active replacement section
    active = config.get("Settings", "ActiveReplacement", fallback="Replacement1")
    if active not in config:
        active = "Replacement1"

    keywords, replacements = [], []
    for k, r in config.items(active):
        keywords.append(k)
        replacements.append(r)

    return keywords, replacements

# ログファイルは exe と同じディレクトリに保存する
log_filename = f"{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}_ReplaceLog.txt"
log_path = os.path.join(base_path, log_filename)

def log(message):
    with open(log_path, 'a', encoding='utf-8') as f:
        f.write(f"{message}\n")

def edit_ini(path):
    if os.path.exists(path):
        config.read(path, encoding="utf-8")

    # Upgrade old style section if necessary
    if "Replacements" in config and "Replacement1" not in config:
        config["Replacement1"] = {}
        for k, v in config["Replacements"].items():
            config["Replacement1"][k] = v
        config.remove_section("Replacements")

    for i in range(1, 6):
        section = f"Replacement{i}"
        if section not in config:
            config[section] = {}

    if "Settings" not in config:
        config["Settings"] = {"ActiveReplacement": "Replacement1"}

    root = tk.Tk()
    root.title("間違いやすい用語チェック.ini 編集")

    instructions = (
        "\n"
        "実行方法：チェックしたいファイルをexeにドロップして実行します。\n\n"
        "設定方法：ドロップせずに実行すると設定モード（この画面）が起動します。\n"
        "　　　　　「選択中の単語グループ」の中の単語が検索対象になります。"
        "\n"

    )
    tk.Label(root, text=instructions, justify="left").pack(padx=10, pady=(10, 0), anchor="w")

    style = ttk.Style()
    style.configure("Bold.TNotebook.Tab", padding=[6, 4])
    style.map(
        "Bold.TNotebook.Tab",
        font=[("selected", ("TkDefaultFont", 9, "bold")), ("!selected", ("TkDefaultFont", 9, "normal"))],
    )

    notebook = ttk.Notebook(root, style="Bold.TNotebook")
    notebook.pack(padx=10, pady=10)

    widgets = {}

    def make_refresh(sec, lb):
        def _refresh():
            lb.delete(0, tk.END)
            for k, v in config[sec].items():
                lb.insert(tk.END, f"{k} = {v}")
        return _refresh

    for i in range(1, 6):
        sec = f"Replacement{i}"
        frame = tk.Frame(notebook)
        notebook.add(frame, text=sec)

        lb = tk.Listbox(frame, width=50, height=15)
        lb.grid(row=0, column=0, columnspan=4, sticky="nsew")
        sb = tk.Scrollbar(frame, orient="vertical", command=lb.yview)
        sb.grid(row=0, column=4, sticky="ns")
        lb.config(yscrollcommand=sb.set)

        tk.Label(frame, text="検索ワード").grid(row=1, column=0, padx=5)
        tk.Label(frame, text="変更推奨ワード").grid(row=1, column=1, padx=5)

        ek = tk.Entry(frame)
        ev = tk.Entry(frame)
        ek.grid(row=2, column=0, padx=5, pady=5)
        ev.grid(row=2, column=1, padx=5, pady=5)

        refresh_func = make_refresh(sec, lb)

        def on_select(event=None, lb=lb, ek=ek, ev=ev):
            if not lb.curselection():
                return
            item = lb.get(lb.curselection()[0])
            k, v = item.split(" = ", 1)
            ek.delete(0, tk.END)
            ek.insert(0, k)
            ev.delete(0, tk.END)
            ev.insert(0, v)

        def on_add_update(event=None, sec=sec, lb=lb, ek=ek, ev=ev):
            k = ek.get().strip()
            v = ev.get().strip()
            if not k or not v:
                return
            config[sec][k] = v

            # Update or append the item in the listbox directly
            found = None
            for i in range(lb.size()):
                item_k = lb.get(i).split(" = ", 1)[0]
                if item_k == k:
                    found = i
                    break
            if found is not None:
                lb.delete(found)
                lb.insert(found, f"{k} = {v}")
            else:
                lb.insert(tk.END, f"{k} = {v}")
                lb.yview_moveto(1)

            lb.update_idletasks()
            ek.delete(0, tk.END)
            ev.delete(0, tk.END)
            ek.focus_set()

        def on_delete(event=None, sec=sec, lb=lb):
            if not lb.curselection():
                return
            index = lb.curselection()[0]
            item = lb.get(index)
            if not messagebox.askyesno("削除確認", "選択した項目を削除しますか？"):
                return
            k = item.split(" = ", 1)[0]
            config[sec].pop(k, None)
            lb.delete(index)
            lb.update_idletasks()

        ek.bind("<Return>", on_add_update)
        ev.bind("<Return>", on_add_update)

        lb.bind("<<ListboxSelect>>", on_select)

        btn_add = tk.Button(frame, text="追加/更新", width=10, command=on_add_update)
        btn_add.grid(row=2, column=2, padx=5, pady=5)
        btn_del = tk.Button(frame, text="削除", width=10, command=on_delete)
        btn_del.grid(row=2, column=3, padx=5, pady=5)

        widgets[sec] = {
            "entry_key": ek,
            "refresh": refresh_func,
            "on_add": on_add_update,
            "on_delete": on_delete,
        }

        refresh_func()

    active_var = tk.StringVar(value=config["Settings"].get("ActiveReplacement", "Replacement1"))
    tk.Label(root, text="選択中の単語グループ").pack()
    selector = ttk.Combobox(root, textvariable=active_var, values=[f"Replacement{i}" for i in range(1, 6)], state="readonly")
    selector.pack(pady=5)

    def current_section():
        idx = notebook.index(notebook.select())
        return f"Replacement{idx + 1}"

    def on_save(event=None):
        config.setdefault("Settings", {})["ActiveReplacement"] = active_var.get()
        with open(path, "w", encoding="utf-8") as f:
            config.write(f)
        root.destroy()

    def on_add(event=None):
        widgets[current_section()]["on_add"]()

    def on_del(event=None):
        widgets[current_section()]["on_delete"]()

    root.bind("<Delete>", on_del)
    root.bind("<Control-s>", on_save)

    def focus_current(event=None):
        widgets[current_section()]["entry_key"].focus_set()

    notebook.bind("<<NotebookTabChanged>>", focus_current)
    root.after(100, focus_current)

    tk.Button(root, text="保存して終了", command=on_save).pack(pady=5)

    root.mainloop()

def search_text_in_docx(path, keywords, replacements):
    from docx import Document
    from docx.oxml.text.paragraph import CT_P
    from docx.text.paragraph import Paragraph

    log(f"――――　ファイル: {os.path.basename(path)}　――――")
    doc = Document(path)
    current_heading = "章番号不明"

    for element in doc.element.body.iter():
        if isinstance(element, CT_P):
            para = Paragraph(element, doc)
            if para.style.name.startswith("Heading"):
                # 見出し段落の先頭番号を見出し番号として扱う（例: "1.1 概要" → "1.1"）
                split_text = para.text.strip().split()
                if split_text and any(char.isdigit() for char in split_text[0]):
                    current_heading = split_text[0]
                else:
                    current_heading = para.text.strip()
            for k, r in zip(keywords, replacements):
                if k in para.text:
                    log(f"{current_heading}: '{k}' → '{r}'")

def search_text_in_xlsx(path, keywords, replacements):
    from openpyxl import load_workbook
    log(f"――――　ファイル: {os.path.basename(path)}　――――")
    wb = load_workbook(path)
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    for k, r in zip(keywords, replacements):
                        if k in cell.value:
                            log(f"シート'{sheet.title}' セル{cell.coordinate}: '{k}' → '{r}'")

def search_text_in_pptx(path, keywords, replacements):
    from pptx import Presentation
    from pptx.table import Table
    log(f"――――　ファイル: {os.path.basename(path)}　――――")
    prs = Presentation(path)

    def walk_shapes(shapes, slide_no):
        for shape in shapes:
            if shape.has_text_frame:
                for k, r in zip(keywords, replacements):
                    if k in shape.text:
                        log(f"スライド{slide_no}: '{k}' → '{r}'")
            if shape.has_table:
                table: Table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for k, r in zip(keywords, replacements):
                            if k in cell.text:
                                log(f"スライド{slide_no}: '{k}' → '{r}'")
            if hasattr(shape, "shapes"):
                walk_shapes(shape.shapes, slide_no)

    for i, slide in enumerate(prs.slides, 1):
        walk_shapes(slide.shapes, i)

def search_text_in_txt(path, keywords, replacements):
    """Search keywords in a plain text file and log matches with line numbers."""
    log(f"――――　ファイル: {os.path.basename(path)}　――――")
    with open(path, encoding="utf-8", errors="ignore") as f:
        for lineno, line in enumerate(f, 1):
            for k, r in zip(keywords, replacements):
                if k in line:
                    log(f"行{lineno}: '{k}' → '{r}'")

def process_files(filepaths):
    keywords, replacements = load_replacements()
    # 処理対象のファイル一覧を冒頭に記録する
    names = ', '.join(os.path.basename(p) for p in filepaths)
    log(f"比較ファイル: {names}")
    log("")
    for path in filepaths:
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext == ".docx":
                search_text_in_docx(path, keywords, replacements)
                log("")
                log("")
            elif ext == ".xlsx":
                search_text_in_xlsx(path, keywords, replacements)
                log("")
                log("")
            elif ext == ".pptx":
                search_text_in_pptx(path, keywords, replacements)
                log("")
                log("")
            elif ext == ".txt":
                search_text_in_txt(path, keywords, replacements)
                log("")
                log("")
            else:
                log(f"[SKIP] 未対応の拡張子: {path}")
        except Exception as e:
            log(f"[ERROR] {path} 処理中にエラー: {str(e)}")

    os.system(f"notepad.exe {log_path}")

if __name__ == "__main__":
    if len(sys.argv) > 1 and os.path.exists(ini_path):
        input_files = sys.argv[1:]
        process_files(input_files)
    else:
        edit_ini(ini_path)
        sys.exit(0)

